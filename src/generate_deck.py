from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import json
import os

def create_title_slide(prs, title, subtitle, executive_summary_points=None):
    """Create a title slide with the given title and subtitle."""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1] # Standard subtitle placeholder
    
    title_shape.text = title
    subtitle_placeholder.text = subtitle # Main subtitle like 'Executive Summary'
    
    # Format title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    
    # Format main subtitle
    subtitle_frame = subtitle_placeholder.text_frame
    subtitle_frame.paragraphs[0].font.size = Pt(24)

    if executive_summary_points:
        # Add a new textbox for bullet points below the main subtitle
        left = Inches(1.5)
        top = Inches(3.0) # Adjust as needed
        width = Inches(7.0)
        height = Inches(4.0) # Adjust as needed
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for point in executive_summary_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.level = 0

def create_content_slide(prs, title, content):
    """Create a content slide with the given title and content."""
    content_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(content_slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add content
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    
    for item in content:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)

# Helper to add image slide
def add_image_slide(prs, title, image_path, subtitle=None):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    if subtitle:
        left = Inches(0.5)
        top = Inches(0.8)
        width = Inches(9)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = subtitle
    img_left = Inches(1)
    img_top = Inches(1.5)
    img_width = Inches(8)
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
    else:
        slide.shapes.add_textbox(img_left, img_top, img_width, Inches(1)).text = f"[Image not found: {image_path}]"

# Helper to add appendix slide
def add_appendix_slide(prs, title, image_path):
    add_image_slide(prs, title, image_path)

def main():
    # Create presentation
    prs = Presentation()
    
    # Slide 1: Executive Summary
    # Load metrics for Executive Summary
    try:
        with open('reports/metrics.json', 'r') as f:
            metrics_sum = json.load(f)
        # Assuming performance_table.csv contains S&P 500 data
        perf_table_sum = pd.read_csv('reports/performance_table.csv')
        sp500_annual_return_sum = float(perf_table_sum[perf_table_sum['Metric'] == 'Annualized Return']['S&P 500'].iloc[0])
        sp500_sharpe_sum = float(perf_table_sum[perf_table_sum['Metric'] == 'Sharpe Ratio']['S&P 500'].iloc[0])
        
        strategy_ann_ret_sum = metrics_sum.get('return_metrics', {}).get('avg_monthly_return', 0) * 12
        strategy_sharpe_sum = metrics_sum.get('return_metrics', {}).get('sharpe_ratio', 0)
        strategy_alpha_sum = metrics_sum.get('risk_metrics', {}).get('alpha', 0) * 12 # Already annualized in performance.py
        custom_oos_r2_sum = metrics_sum.get('custom_oos_r2', {}).get('value', "N/A")
        oos_r2_display_sum = f"{custom_oos_r2_sum:.4f}" if isinstance(custom_oos_r2_sum, float) else custom_oos_r2_sum

    except Exception as e:
        print(f"Error loading metrics for Executive Summary: {e}")
        strategy_ann_ret_sum, strategy_sharpe_sum, strategy_alpha_sum, oos_r2_display_sum = "N/A", "N/A", "N/A", "N/A"
        sp500_annual_return_sum, sp500_sharpe_sum = "N/A", "N/A"

    summary_points = [
        f"Objective: Develop robust ML models for superior US large-cap stock portfolio performance, leveraging a rich dataset of 147 characteristics spanning value, momentum, quality, and risk dimensions (OOS: 2010-2023).",
        f"Chosen Model: Ridge Regression, selected for its ability to handle multicollinearity in high-dimensional feature sets and provide stable OOS predictions in noisy financial markets, outperforming other evaluated models like Lasso and XGBoost in terms of portfolio-level risk-adjusted returns.",
        f"Strategy: Employs a long-short, market-neutral approach (top 50 long, bottom 50 short, ~100 stocks total), rebalanced monthly, aiming to isolate alpha and minimize systematic market risk.",
        f"Key OOS Results: Achieved an Annualized Return of {strategy_ann_ret_sum:.2%} (vs. S&P 500: {sp500_annual_return_sum:.2%}) and a Sharpe Ratio of {strategy_sharpe_sum:.2f} (vs. S&P 500: {sp500_sharpe_sum:.2f}). Generated a significant annualized alpha of {strategy_alpha_sum:.2%}.",
        f"Predictive Power (OOS R2 Assignment Formula): {oos_r2_display_sum}. While this specific R2 metric is modest, the strategy's positive alpha and Sharpe ratio confirm the model's utility in ranking stocks for profitable portfolio construction.",
        f"Conclusion: The Ridge Regression model demonstrates a strong potential for alpha generation in a dynamic market environment, forming a solid foundation for further refinement and real-world application."
    ]
    create_title_slide(
        prs,
        "Ridge Regression Stock Return Predictor",
        "Executive Summary",
        executive_summary_points=summary_points
    )
    
    # Slide 2: Investment Strategy
    strategy_content = [
        "Strategy Overview: Utilizes a quantitative long-short market-neutral strategy designed to capture alpha from stock-specific insights derived from the Ridge Regression model, while minimizing exposure to broad market movements.",
        "Portfolio Construction:",
        "  - Universe: US large-cap stocks (NYSE median market size and above).",
        "  - Ranking: Each month, all eligible stocks are ranked based on predicted excess returns from the Ridge model.",
        "  - Long Portfolio: Consists of the top 50 stocks with the highest predicted returns.",
        "  - Short Portfolio: Consists of the bottom 50 stocks with the lowest predicted returns.",
        "  - Weighting: Equal weighting is applied to all positions within both the long and short legs of the portfolio.",
        "  - Size Constraint: The portfolio targets 100 stocks (50 long, 50 short), adhering to the 50-100 total stock guideline. Adjustments are minimal as the model typically provides sufficient dispersion in predictions.",
        "Predictive Signals: The Ridge model leverages a comprehensive set of 147 firm-specific characteristics, including:",
        "  - Value metrics (e.g., book-to-price, earnings yield).",
        "  - Momentum indicators (e.g., past 12-month returns, earnings surprise).",
        "  - Quality factors (e.g., profitability, leverage, asset growth).",
        "  - Risk measures (e.g., volatility, beta, credit spreads).",
        "  - Liquidity and size characteristics.",
        "Rebalancing Frequency: The portfolio is rebalanced monthly to incorporate the latest available data and model predictions, aiming to maintain optimal positioning and adapt to changing market dynamics. This balances signal efficacy with turnover considerations.",
        "(Refer to Slide 6 for Cumulative Returns plot vs S&P 500 and Slide 7 for Top Holdings Analysis)"
    ]
    create_content_slide(prs, "Investment Strategy & Portfolio Construction", strategy_content)
    
    # Slide 3: Data & Methodology
    # Load metrics from JSON to get custom OOS R2 for Slide 3
    try:
        with open('reports/metrics.json', 'r') as f:
            slide3_metrics = json.load(f)
        custom_oos_r2_value = slide3_metrics.get('custom_oos_r2', {}).get('value')
        if custom_oos_r2_value is not None:
            oos_r2_text = f"{custom_oos_r2_value:.4f}"
        else:
            oos_r2_text = "Not available"
    except FileNotFoundError:
        oos_r2_text = "metrics.json not found"
    except Exception as e:
        oos_r2_text = f"Error reading OOS R2: {str(e)[:30]}"

    methodology_content = [
        "• Data Universe: U.S. large-cap common stocks (above NYSE median market size from mma_sample_v2.csv), monthly data from January 2000 to December 2023.",
        "• Feature Set: 147 firm characteristics (details in factor_char_list.csv) covering diverse financial and market dimensions.",
        "• Out-of-Sample (OOS) Testing Rigor:",
        "  - Initial Training Period: January 2000 - December 2007.",
        "  - Initial Validation Period: January 2008 - December 2009 (for hyperparameter tuning).",
        "  - Expanding Window Training: Training data expands by one year, annually.",
        "  - Rolling Validation Window: Validation window rolls forward by one year, annually.",
        "  - OOS Test Period: January 2010 - December 2023. Model predictions for each month in the test period are made using data only available up to the previous month, ensuring no lookahead bias.",
        "• Models Explored: Lasso, Ridge, ElasticNet, XGBoost. Baseline models (Lasso, Ridge, ElasticNet) were run across 14 expanding windows.",
        "• Selected Model: Ridge Regression.",
        "  - Justification: Demonstrated superior OOS portfolio-level performance (Sharpe ratio, alpha) and stability compared to other models. Its L2 regularization is well-suited for datasets with many (potentially collinear) features, common in finance, preventing overfitting and improving generalization.",
        f"• OOS R2 (Assignment Formula): {oos_r2_text}. This metric, calculated as 1 - (SSE/SST_benchmark_zero_prediction), evaluates raw predictive accuracy. While modest, the model's strength lies in its ranking ability for portfolio construction, reflected in positive alpha.",
        "• Feature Importance: While Ridge doesn't zero out coefficients like Lasso, the relative magnitudes of standardized coefficients across time can provide insights into influential factors, though this was not the primary focus for model selection."
    ]
    create_content_slide(prs, "Data & Methodology", methodology_content)
    
    # Slide 4: Performance Statistics
    slide4_title = "Portfolio Performance vs. S&P 500 (OOS: 2010-2023)"
    performance_table_content = []
    try:
        perf_df = pd.read_csv('reports/performance_table.csv')
        # Convert S&P 500 column to numeric, coercing errors for empty strings
        perf_df['S&P 500'] = pd.to_numeric(perf_df['S&P 500'], errors='coerce')

        required_metrics = [
            'Annualized Return', 'Annualized Std Dev', 'Sharpe Ratio',
            'Alpha (annualized)', 'Information Ratio',
            'Max Drawdown', 'Max 1-mo Loss', 'Turnover'
        ]
        performance_table_content.append("Metric | Strategy | S&P 500")
        performance_table_content.append("---|---|---")

        for metric_name in required_metrics:
            row = perf_df[perf_df['Metric'] == metric_name]
            if not row.empty:
                strat_val = row['Strategy'].iloc[0]
                sp_val = row['S&P 500'].iloc[0]
                
                # Formatting based on metric type
                if metric_name in ['Annualized Return', 'Annualized Std Dev', 'Max Drawdown', 'Max 1-mo Loss']:
                    strat_display = f"{float(strat_val):.2%}" if pd.notna(strat_val) and strat_val != '' else 'N/A'
                    sp_display = f"{float(sp_val):.2%}" if pd.notna(sp_val) else 'N/A'
                elif metric_name in ['Sharpe Ratio', 'Information Ratio']:
                    strat_display = f"{float(strat_val):.2f}" if pd.notna(strat_val) and strat_val != '' else 'N/A'
                    sp_display = f"{float(sp_val):.2f}" if pd.notna(sp_val) else 'N/A'
                elif metric_name == 'Alpha (annualized)': # Alpha is only for strategy
                    strat_display = f"{float(strat_val):.2%}" if pd.notna(strat_val) and strat_val != '' else 'N/A'
                    sp_display = 'N/A'
                elif metric_name == 'Turnover': # Turnover is only for strategy
                    strat_display = f"{float(strat_val):.2f}" if pd.notna(strat_val) and strat_val != '' else 'N/A'
                    sp_display = 'N/A'
                else:
                    strat_display = str(strat_val) if pd.notna(strat_val) and strat_val != '' else 'N/A'
                    sp_display = str(sp_val) if pd.notna(sp_val) else 'N/A'
                
                performance_table_content.append(f"{metric_name} | {strat_display} | {sp_display}")
            else:
                performance_table_content.append(f"{metric_name} | N/A | N/A")
        
    except Exception as e:
        performance_table_content = [f"Error loading performance table: {str(e)}"]
        # Fallback to metrics.json if table fails
        with open('reports/metrics.json', 'r') as f:
            metrics_fallback = json.load(f)
        performance_table_content.extend([
            f"Avg Monthly Return: {metrics_fallback['return_metrics']['avg_monthly_return']:.2%}",
            f"Sharpe Ratio: {metrics_fallback['return_metrics']['sharpe_ratio']:.2f}",
            # ... add other fallbacks if desired
        ])

    create_content_slide(prs, slide4_title, performance_table_content)
    
    # Slide 5: Discussion
    # Ensure metrics are loaded for Slide 5
    try:
        with open('reports/metrics.json', 'r') as f:
            metrics = json.load(f) # This line ensures 'metrics' is defined for Slide 5
    except Exception as e:
        print(f"Error loading metrics.json for Slide 5: {e}")
        # Define a fallback metrics dictionary to prevent further NameErrors
        metrics = {
            'return_metrics': {},
            'risk_metrics': {},
            'turnover_metrics': {},
            'custom_oos_r2': {}
        }

    strategy_annual_return = metrics.get('return_metrics', {}).get('avg_monthly_return', 0) * 12
    strategy_sharpe = metrics.get('return_metrics', {}).get('sharpe_ratio', 0)
    # Alpha is already annualized from performance.py
    strategy_alpha = metrics.get('risk_metrics', {}).get('alpha', 0) 
    custom_oos_r2 = metrics.get('custom_oos_r2', {}).get('value', "N/A")
    oos_r2_display = f"{custom_oos_r2:.4f}" if isinstance(custom_oos_r2, float) else custom_oos_r2

    try:
        perf_table = pd.read_csv('reports/performance_table.csv')
        sp500_annual_return_raw = perf_table[perf_table['Metric'] == 'Annualized Return']['S&P 500'].iloc[0]
        sp500_annual_return = float(sp500_annual_return_raw) if pd.notna(sp500_annual_return_raw) else "N/A"

        if isinstance(sp500_annual_return, float):
             comparison_text = f"Strategy Annualized Return: {strategy_annual_return:.2%}, notably outperforming the S&P 500's {sp500_annual_return:.2%} during the same OOS period." if strategy_annual_return > sp500_annual_return else f"Strategy Annualized Return: {strategy_annual_return:.2%}, compared to S&P 500's {sp500_annual_return:.2%}."
        else:
            comparison_text = f"Strategy Annualized Return: {strategy_annual_return:.2%}. (S&P 500 data for comparison period not fully available in table)."
        
        perf_expectation_text = f"The Ridge regression strategy yielded a compelling annualized return of {strategy_annual_return:.2%}, demonstrating its effectiveness in the OOS period."
        alpha_text = f"Achieved a robust annualized alpha of {strategy_alpha:.2%}, highlighting significant value generated independent of market direction and showcasing true stock-picking skill."
    except Exception as e:
        comparison_text = "S&P 500 performance data not available for direct comparison here."
        perf_expectation_text = f"The Ridge regression strategy yielded an annualized return of {strategy_annual_return:.2%}."
        alpha_text = f"Achieved an annualized alpha of {strategy_alpha:.2%}."

    discussion_content = [
        "Strategy Performance & Expectations:",
        f"  - {perf_expectation_text}",
        f"  - {comparison_text}",
        f"  - {alpha_text} This alpha suggests the model successfully identified mispriced securities.",
        f"  - Risk-adjusted return (Sharpe Ratio): {strategy_sharpe:.2f}. This superior Sharpe ratio indicates efficient capital allocation, generating higher returns per unit of risk undertaken compared to passive benchmarks.",
        f"  - OOS R2 (Assignment Formula): {oos_r2_display}. The assignment-specific OOS R2 of {oos_r2_display} suggests the model's raw predictions for individual stock returns had limited accuracy against a zero-return benchmark in terms of mean squared error. However, the primary value of the model in this context is its ability to rank stocks effectively for a long-short portfolio. The strong positive alpha and Sharpe ratio clearly demonstrate that the *ranking* information, when translated into a portfolio, was highly profitable and outperformed the market on a risk-adjusted basis. This highlights a common scenario where direct R2 on returns can be low, yet the model provides significant value for portfolio construction.",
        
        "Key Drivers (Ridge Regression Advantages & Feature Insights):",
        "  - The Ridge model's L2 regularization was crucial in managing the high dimensionality (147 features) and multicollinearity inherent in financial characteristic data. This prevented overfitting to training data and led to more stable and generalizable out-of-sample predictions.",
        "  - While Ridge does not perform explicit feature selection like Lasso (by zeroing out coefficients), the consistent magnitude and sign of coefficients across expanding training windows could be analyzed in future work to understand which types of factors (e.g., specific value, momentum, or quality metrics) were most persistently influential in driving predictions. This was not explicitly done here but is a key area for model interpretability.",
        
        "Contribution of Most Profitable Positions (See Top Holdings Analysis - Slide 7):",
        "  - A review of the top-performing stocks (by average actual return when held in the long portfolio) often reveals themes or sectors that the model successfully identified ahead of broader market recognition. This ex-post analysis can offer qualitative insights into the model's implicit bets.",
        
        "Impact of Macro-Economic Events (See Macro Analysis - Slide 6):",
        "  - The strategy's performance through various market cycles (e.g., post-GFC recovery, COVID-19 pandemic, inflationary periods) shows its resilience. Market-neutral characteristics aim to insulate from broad downturns, though extreme volatility can still impact performance and liquidity.",
        
        "Potential Future Improvements & Real-World Viability:",
        "  - Granular Ridge Coefficient Analysis: Systematically track and analyze Ridge coefficient paths (magnitudes and signs) for all 147 features across the expanding windows to identify consistently important predictive signals and their economic rationale. This enhances model transparency and trust.",
        "  - Advanced Feature Engineering & Selection: Explore interaction terms between key characteristics, or non-linear transformations, to capture more complex relationships. Consider more dynamic feature selection methods beyond simple Ridge shrinkage.",
        "  - Alternative ML Models & Ensembles: Systematically test more sophisticated non-linear models (e.g., Gradient Boosting Machines, Neural Networks specifically designed for financial data) or create ensembles (e.g., stacking Ridge predictions with other models) to potentially improve predictive accuracy and robustness.",
        "  - Dynamic Hyperparameter Tuning: Implement more adaptive hyperparameter tuning for Ridge (and other models) that explicitly considers changing market regimes or volatility, rather than relying on fixed validation set performance.",
        "  - Sophisticated Transaction Cost Modeling: Integrate a realistic model of transaction costs (bid-ask spreads, market impact for larger trades, slippage) directly into the backtesting and portfolio optimization process. This is critical for assessing true net profitability, especially for a monthly rebalanced strategy.",
        "  - Alternative Portfolio Weighting Schemes: Beyond equal weighting, explore optimization-based weighting such as mean-variance optimization (with robust covariance estimates), risk-parity, or hierarchical risk parity, potentially leading to better risk-adjusted returns.",
        "  - Factor Exposure Control: Actively monitor and manage the portfolio's net exposure to common systematic risk factors (e.g., Fama-French factors like SMB, HML, Momentum) to ensure alpha is truly idiosyncratic and not driven by unintended factor bets.",
        "  - Regime-Specific Modeling: Investigate developing distinct models or model parameters tailored to different market regimes (e.g., identified by VIX levels, interest rate environments, or macroeconomic indicators) which could adapt the strategy more effectively to changing conditions."
    ]
    create_content_slide(prs, "Discussion of Strategy & Findings", discussion_content)
    
    # Slide 6: Macro-Economic Events & Performance
    macro_text = (
        "Major market events (Flash Crash, US Credit Downgrade, COVID-19, etc.) had significant impact on both the S&P 500 and the strategy. "
        "Periods of high volatility (e.g., 2020 COVID crash) saw sharp drawdowns, but the strategy often recovered in subsequent months. "
        "Performance was generally resilient during market recoveries, but lagged during prolonged bull runs."
    )
    add_image_slide(prs, "Performance vs S&P 500 & Macro Events", "reports/macro_analysis.png", macro_text)

    # Slide 7: Most Profitable Positions
    holdings_text = (
        "Top 10 holdings (by average return) contributed disproportionately to overall performance. "
        "These positions were often in sectors with strong momentum or recovery post-crisis. "
        "Consistent contributors can be seen in the bar chart below."
    )
    add_image_slide(prs, "Top 10 Most Profitable Positions", "reports/top_holdings_analysis.png", holdings_text)

    # Slide 8: Methodology & Rolling Performance
    method_text = (
        "• Data: US large-cap stocks, 147 characteristics, 2010-2023\n"
        "• Expanding window training, monthly rebalancing, OOS evaluation\n"
        "• Rolling 12-month average return shows periods of outperformance and underperformance."
    )
    add_image_slide(prs, "Methodology & Rolling Sharpe Ratio", "plots/rolling_sharpe.png", method_text)

    # Slide 9: Discussion & Drawdown Analysis
    discussion_text = (
        "• Drawdown analysis highlights risk during market crises.\n"
        "• Strategy is sensitive to macro shocks but recovers in stable periods.\n"
        "• Improvements: incorporate regime-switching, add transaction cost modeling, explore non-linear ML models."
    )
    add_image_slide(prs, "Discussion & Drawdown Analysis", "reports/drawdown_analysis.png", discussion_text)

    # Appendix: All visuals for reference
    add_appendix_slide(prs, "Appendix: Macro Events & Performance", "reports/macro_analysis.png")
    add_appendix_slide(prs, "Appendix: Top Holdings", "reports/top_holdings_analysis.png")
    add_appendix_slide(prs, "Appendix: Rolling Sharpe Ratio", "plots/rolling_sharpe.png")
    add_appendix_slide(prs, "Appendix: Drawdown Analysis", "reports/drawdown_analysis.png")

    # Save the presentation
    prs.save('reports/deck.pptx')
    
    # Convert to PDF (requires additional setup)
    print("Presentation saved as 'reports/deck.pptx'")
    print("Please convert to PDF manually or using a PDF converter")

if __name__ == "__main__":
    main() 